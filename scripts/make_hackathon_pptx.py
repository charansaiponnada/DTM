"""Generate Hackathon Final Round PPTX with beige theme."""
import json, os, textwrap
from pathlib import Path

BASE = Path("C:/projects/DTM")
OUTPUT_DIR = BASE / "data" / "output"

# ── Colour Palette (Beige Theme) ──────────────────────────────────────
CREAM   = (0xF5, 0xF0, 0xE8)   # slide background
TAN     = (0xE8, 0xDC, 0xC8)   # section banner
NAVY    = (0x2C, 0x3E, 0x50)   # headings, accent bar
BROWN   = (0x8B, 0x73, 0x55)   # secondary accent
SIENNA  = (0xA0, 0x52, 0x2D)   # highlight
DARK    = (0x3D, 0x2B, 0x1F)   # body text
WHITE   = (0xFF, 0xFF, 0xFF)
LIGHT   = (0xF0, 0xEB, 0xE1)
GREY    = (0x7F, 0x8C, 0x8D)
ORANGE  = (0xD3, 0x54, 0x00)

def load_metrics(village):
    p = OUTPUT_DIR / village / "metrics.json"
    return json.loads(p.read_text()) if p.exists() else {}

devdi  = load_metrics("DEVDI")
khapreta = load_metrics("KHAPRETA")
villages_data = [
    {
        "name": "DEVDI",
        "file": "DEVDI_511671.las",
        "pts": "64.6 M",
        "area": "0.81 km\u00b2",
        "density": "65 pts/m\u00b2",
        "tiles": 6,
        "elev_range": "-31.2 to -7.9 m",
        "elev_mean": -18.3,
        "relief": 23.3,
        **devdi,
    },
    {
        "name": "KHAPRETA",
        "file": "KHAPRETA_510206.laz",
        "pts": "163.7 M",
        "area": "0.67 km\u00b2",
        "density": "245 pts/m\u00b2",
        "tiles": 4,
        "elev_range": "115.3 to 151.0 m",
        "elev_mean": 133.8,
        "relief": 35.8,
        **khapreta,
    },
]

def rgb(color):
    from pptx.dml.color import RGBColor
    return RGBColor(*color)

def make_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    W = prs.slide_width
    H = prs.slide_height

    def add_slide():
        layout = prs.slide_layouts[6]
        return prs.slides.add_slide(layout)

    def add_bg(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb(color)

    def shape(slide, l, t, w, h, fill=None, line=None, radius=None):
        if radius:
            sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
            sh.adjustments[0] = radius
        else:
            sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        sh.line.fill.background()
        if fill:
            sh.fill.solid()
            sh.fill.fore_color.rgb = rgb(fill)
        else:
            sh.fill.background()
        if line:
            sh.line.color.rgb = rgb(line)
            sh.line.width = Pt(1.5)
        return sh

    def _add(slide, l, t, w, h, text, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, font="Calibri", italic=False, anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        try:
            tf.vertical_anchor = anchor
        except:
            pass
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = rgb(color)
        p.font.name = font
        p.alignment = align
        return tb

    def add_bullets(slide, l, t, w, h, items, size=16, color=DARK, font="Calibri",
                    spacing=6, symbol=None):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            prefix = f"{symbol} " if symbol else "\u2022 "
            p.text = f"{prefix}{item}"
            p.font.size = Pt(size)
            p.font.color.rgb = rgb(color)
            p.font.name = font
            p.space_after = Pt(spacing)
            p.level = 0
        return tb

    def section_header(slide, title, subtitle=None):
        shape(slide, 0, 0, W, Inches(1.3), fill=NAVY)
        _add(slide, Inches(0.7), Inches(0.2), Inches(10), Inches(0.7),
             title, size=30, bold=True, color=WHITE)
        shape(slide, Inches(0.7), Inches(1.1), Inches(2.5), Inches(0.06), fill=SIENNA)
        if subtitle:
            _add(slide, Inches(0.7), Inches(1.4), Inches(11), Inches(0.4),
                 subtitle, size=13, color=GREY)

    def make_table(slide, l, t, w, h, rows, cols, data, col_widths=None,
                   header_color=NAVY, alt_color=LIGHT):
        tbl = slide.shapes.add_table(rows, cols, l, t, w, h).table
        for r in range(rows):
            for c in range(cols):
                cell = tbl.cell(r, c)
                cell.text = str(data[r][c]) if r < len(data) and c < len(data[r]) else ""
                for par in cell.text_frame.paragraphs:
                    par.font.size = Pt(12)
                    par.font.name = "Calibri"
                    if r == 0:
                        par.font.bold = True
                        par.font.color.rgb = rgb(WHITE)
                        par.alignment = PP_ALIGN.CENTER
                    else:
                        par.font.color.rgb = rgb(DARK)
                if r == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = rgb(header_color)
                elif r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = rgb(alt_color)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = rgb(WHITE)
        if col_widths:
            for i, w2 in enumerate(col_widths):
                tbl.columns[i].width = w2
        return tbl

    # ═════════════════════════════════════════════════════════════════
    #  SLIDE 1 — TITLE
    # ═════════════════════════════════════════════════════════════════
    s = add_slide()
    add_bg(s, CREAM)
    # Top accent bar
    shape(s, 0, 0, W, Inches(0.08), fill=SIENNA)
    # Left decorative block
    shape(s, 0, Inches(1.5), Inches(0.5), Inches(4.5), fill=NAVY)
    # Main title
    _add(s, Inches(1.2), Inches(1.8), Inches(11), Inches(1.2),
         "DTM Drainage AI Pipeline", size=48, bold=True, color=NAVY)
    _add(s, Inches(1.2), Inches(3.0), Inches(11), Inches(0.8),
         "Automated DTM Generation \u00b7 Waterlogging Prediction \u00b7 Drainage Network Design",
         size=20, color=BROWN)
    # Separator
    shape(s, Inches(1.2), Inches(3.9), Inches(4), Inches(0.04), fill=SIENNA)
    # Subtitle
    _add(s, Inches(1.2), Inches(4.2), Inches(11), Inches(0.5),
         "from Airborne LiDAR Point Clouds under the SVAMITVA Scheme", size=16, color=GREY)
    # Hackathon info
    _add(s, Inches(1.2), Inches(5.0), Inches(11), Inches(0.4),
         "MoPR Geospatial Intelligence Hackathon \u2014 Problem Statement 2", size=16, bold=True, color=SIENNA)
    # Villages
    _add(s, Inches(1.2), Inches(5.6), Inches(11), Inches(0.4),
         "Villages: DEVDI (Gujarat)  \u00b7  KHAPRETA (Gujarat)", size=14, color=GREY)
    # Bottom right label
    _add(s, Inches(8), Inches(6.8), Inches(5), Inches(0.5),
         "Siddhartha Academy of Higher Education", size=13, color=BROWN, align=PP_ALIGN.RIGHT)

    # ═════════════════════════════════════════════════════════════════
    #  SLIDE 2 — PROBLEM STATEMENT
    # ═════════════════════════════════════════════════════════════════
    s = add_slide()
    add_bg(s, CREAM)
    section_header(s, "Problem Statement")

    items = [
        "Ministry of Panchayati Raj (MoPR) SVAMITVA scheme captures drone LiDAR of ~10 village abadi areas in Gujarat",
        "Challenge: Generate high-resolution Digital Terrain Models (DTMs) from raw LiDAR point clouds",
        "Identify waterlogging hotspot zones using AI/ML on terrain-derived features",
        "Design a cost-optimal drainage network for each village",
        "Deliver all outputs in OGC-compliant GIS formats (COG raster + GPKG vector)",
        "Target: Village-scale solutions for flood mitigation and infrastructure planning",
    ]
    add_bullets(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.5), items, size=17, color=DARK)

    # Key requirements box
    shape(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.5), fill=TAN, radius=0.05)
    reqs = [
        "Requirement 1: DTM at \u22640.5 m resolution  |  Requirement 2: Waterlogging probability map",
        "Requirement 3: Engineered drainage network with cost estimates  |  Requirement 4: All outputs in open GIS formats",
    ]
    add_bullets(s, Inches(1.0), Inches(5.6), Inches(11.3), Inches(1.3), reqs, size=14, color=NAVY, symbol="\u25b6")

    # ═════════════════════════════════════════════════════════════════
    #  SLIDE 3 — INTRODUCTION
    # ═════════════════════════════════════════════════════════════════
    s = add_slide()
    add_bg(s, CREAM)
    section_header(s, "Introduction & Background")

    # Left column - context
    shape(s, Inches(0.7), Inches(1.7), Inches(5.8), Inches(5.3), fill=WHITE, radius=0.05)
    intro_items = [
        "Gujarat faces recurrent waterlogging during monsoon; accurate DTMs are critical for drainage planning",
        "SVAMITVA (Survey of Villages and Mapping with Improvised Technology in Village Areas) scheme provides high-density LiDAR (65\u2013245 pts/m\u00b2)",
        "Traditional manual drainage design is slow, inconsistent, and does not scale to 10+ villages",
        "Our pipeline automates the entire workflow: raw point cloud \u2192 engineered drainage network in <10 minutes",
    ]
    _add(s, Inches(0.9), Inches(1.9), Inches(5.4), Inches(0.4),
         "Context", size=18, bold=True, color=NAVY)
    add_bullets(s, Inches(0.9), Inches(2.4), Inches(5.4), Inches(4.2), intro_items, size=14, color=DARK)

    # Right column - data stats
    shape(s, Inches(6.8), Inches(1.7), Inches(5.8), Inches(5.3), fill=WHITE, radius=0.05)
    _add(s, Inches(7.0), Inches(1.9), Inches(5.4), Inches(0.4),
         "Input Data", size=18, bold=True, color=NAVY)
    data_info = [
        "DEVDI: 64.6 M points, 0.81 km\u00b2, 65 pts/m\u00b2",
        "KHAPRETA: 163.7 M points, 0.67 km\u00b2, 245 pts/m\u00b2",
        "CRS: EPSG:32643 (WGS 84 / UTM Zone 43N)",
        "Format: LAS 1.4 / LAZ (compressed)",
        "Elevation: DEVDI -31.2 to -7.9 m (below MSL)",
        "Elevation: KHAPRETA 115.3 to 151.0 m",
    ]
    add_bullets(s, Inches(7.0), Inches(2.4), Inches(5.4), Inches(4.2), data_info, size=14, color=DARK)

    # ═════════════════════════════════════════════════════════════════
    #  SLIDE 4 — RELATED WORKS
    # ═════════════════════════════════════════════════════════════════
    s = add_slide()
    add_bg(s, CREAM)
    section_header(s, "Related Works & References")

    refs = [
        ("SMRF Filter", "Mongus, D., & \u017dalik, B. (2012). Parameter-free ground filtering of LiDAR data for automatic DTM generation. ISPRS Journal of Photogrammetry and Remote Sensing."),
        ("PointNet", "Qi, C. R., et al. (2017). PointNet: Deep Learning on Point Sets for 3D Classification and Segmentation. CVPR."),
        ("XGBoost", "Chen, T., & Guestrin, C. (2016). XGBoost: A Scalable Tree Boosting System. KDD."),
        ("pysheds", "Bartos, M. (2020). pysheds: watershed and flow routing in Python. Journal of Open Source Software."),
        ("Manning\u2019s Eq.", "Manning, R. (1891). On the flow of water in open channels and pipes. Transactions of ICE Ireland."),
        ("IS 10430", "Bureau of Indian Standards. IS 10430: Criteria for design of lined canals and guidance for selection of type of lining."),
    ]
    y = Inches(1.7)
    for title, desc in refs:
        shape(s, Inches(0.7), y, Inches(1.8), Inches(0.7), fill=NAVY, radius=0.05)
        _add(s, Inches(0.8), y + Inches(0.08), Inches(1.6), Inches(0.5), title, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add(s, Inches(2.7), y + Inches(0.05), Inches(10), Inches(0.5), desc, size=12, color=DARK)
        y += Inches(0.85)

    # ═════════════════════════════════════════════════════════════════
    #  SLIDE 5 — PIPELINE OVERVIEW
    # ═════════════════════════════════════════════════════════════════
    s = add_slide()
    add_bg(s, CREAM)
    section_header(s, "System Architecture \u2014 6-Stage Pipeline")

    stages = [
        ("1", "Data\nInspection", "LAS/LAZ metadata\ntiling, CRS validation\ndensity & bounds"),
        ("2", "Ground\nClassification", "SMRF filter (PDAL)\nCSF fallback\nRF refinement"),
        ("3", "DTM\nGeneration", "IDW interpolation\n0.5 m resolution\nTerrain derivatives"),
        ("4", "Hydrological\nAnalysis", "Depression filling\nD8 flow routing\nStream extraction"),
        ("5", "Waterlogging\nPrediction", "10-feature stack\nXGBoost classifier\n5-fold CV"),
        ("6", "Drainage\nDesign", "MST optimisation\nManning\u2019s sizing\nCost estimation"),
    ]
    box_w = Inches(1.85)
    box_h = Inches(3.8)
    gap = Inches(0.2)
    start_x = Inches(0.5)
    y_top = Inches(1.8)

    for i, (num, title, desc) in enumerate(stages):
        x = start_x + i * (box_w + gap)
        # Stage number circle
        shape(s, x + Inches(0.65), y_top, Inches(0.55), Inches(0.55), fill=SIENNA, radius=0.3)
        _add(s, x + Inches(0.65), y_top + Inches(0.02), Inches(0.55), Inches(0.5),
             num, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Card body
        shape(s, x, y_top + Inches(0.7), box_w, Inches(3.0), fill=WHITE, line=NAVY)
        _add(s, x + Inches(0.1), y_top + Inches(0.8), box_w - Inches(0.2), Inches(0.8),
             title, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        _add(s, x + Inches(0.1), y_top + Inches(1.7), box_w - Inches(0.2), Inches(1.8),
             desc, size=11, color=GREY, align=PP_ALIGN.CENTER)
        # Arrow between cards
        if i < len(stages) - 1:
            _add(s, x + box_w + Inches(0.02), y_top + Inches(0.2), Inches(0.2), Inches(0.3),
                 "\u25b6", size=14, color=SIENNA, align=PP_ALIGN.CENTER)

    # Input / Output labels
    shape(s, Inches(0.5), Inches(5.8), Inches(2.0), Inches(0.5), fill=NAVY, radius=0.05)
    _add(s, Inches(0.5), Inches(5.82), Inches(2.0), Inches(0.45),
         "INPUT: LAS/LAZ", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    shape(s, Inches(10.8), Inches(5.8), Inches(2.0), Inches(0.5), fill=BROWN, radius=0.05)
    _add(s, Inches(10.8), Inches(5.82), Inches(2.0), Inches(0.45),
         "OUTPUT: COG + GPKG", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Arrow between
    _add(s, Inches(2.6), Inches(5.85), Inches(8.1), Inches(0.4),
         "\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u2500\u25b6",
         size=14, color=NAVY, align=PP_ALIGN.CENTER)

    # ═════════════════════════════════════════════════════════════════
    #  SLIDE 6 — METHODOLOGY Stages 1-3
    # ═════════════════════════════════════════════════════════════════
    s = add_slide()
    add_bg(s, CREAM)
    section_header(s, "Methodology: Data Preparation & Terrain Modelling (Stages 1\u20133)")

    # Stage 1
    shape(s, Inches(0.6), Inches(1.7), Inches(3.8), Inches(2.5), fill=WHITE, line=NAVY)
    shape(s, Inches(0.6), Inches(1.7), Inches(3.8), Inches(0.45), fill=NAVY)
    _add(s, Inches(0.7), Inches(1.72), Inches(3.6), Inches(0.4),
         "Stage 1: Data Inspection", size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    s1_items = [
        "Read LAS/LAZ headers (lazrs/laspy)",
        "Validate CRS (EPSG:32643)",
        "Compute density, bounds, point count",
        "Tile large files: 500m tiles + 25m buffer",
    ]
    add_bullets(s, Inches(0.7), Inches(2.3), Inches(3.6), Inches(1.8), s1_items, size=12, color=DARK, spacing=4)

    # Stage 2
    shape(s, Inches(4.65), Inches(1.7), Inches(3.8), Inches(2.5), fill=WHITE, line=NAVY)
    shape(s, Inches(4.65), Inches(1.7), Inches(3.8), Inches(0.45), fill=NAVY)
    _add(s, Inches(4.75), Inches(1.72), Inches(3.6), Inches(0.4),
         "Stage 2: Ground Classification", size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    s2_items = [
        "SMRF via PDAL (pdal.exe subprocess)",
        "slope=0.15, window=18m, threshold=0.5m",
        "CSF fallback if PDAL unavailable",
        "Random Forest refinement on geometric features",
        "Output: classified_ground.las (ASPRS class 2)",
    ]
    add_bullets(s, Inches(4.75), Inches(2.3), Inches(3.6), Inches(1.8), s2_items, size=12, color=DARK, spacing=4)

    # Stage 3
    shape(s, Inches(8.7), Inches(1.7), Inches(3.8), Inches(2.5), fill=WHITE, line=NAVY)
    shape(s, Inches(8.7), Inches(1.7), Inches(3.8), Inches(0.45), fill=NAVY)
    _add(s, Inches(8.8), Inches(1.72), Inches(3.6), Inches(0.4),
         "Stage 3: DTM Generation", size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    s3_items = [
        "IDW interpolation (k=16, power=2)",
        "cKDTree batch query: 40\u00d7 speedup",
        "0.5 m resolution Cloud-Optimized GeoTIFF",
        "8 terrain derivatives: slope, aspect,",
        "  curvature, TPI, roughness, hillshade",
    ]
    add_bullets(s, Inches(8.8), Inches(2.3), Inches(3.6), Inches(1.8), s3_items, size=12, color=DARK, spacing=4)

    # Key technique highlight
    shape(s, Inches(0.6), Inches(4.5), Inches(11.9), Inches(2.5), fill=TAN, radius=0.05)
    _add(s, Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.4),
         "\U0001f52e Key Technique: IDW Vectorised Interpolation", size=16, bold=True, color=NAVY)
    _add(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.3),
         "Build KD-tree once on ground points \u2192 Query all grid cells in one batched call (workers=-1) \u2192 Vectorised weight computation - no Python loops",
         size=13, color=DARK)
    code_text = (
        "dists, idxs = tree.query(grid_points, k=16, workers=-1)\n"
        "weights = 1.0 / np.maximum(dists**2, 1e-10)\n"
        "z_interp = (weights * ground_z[idxs]).sum(axis=1) / weights.sum(axis=1)"
    )
    shape(s, Inches(0.8), Inches(5.5), Inches(11.3), Inches(1.2), fill=NAVY, radius=0.04)
    _add(s, Inches(1.0), Inches(5.55), Inches(10.9), Inches(1.1),
         code_text, size=11, color=WHITE, font="Consolas")

    # ═════════════════════════════════════════════════════════════════
    #  SLIDE 7 — METHODOLOGY Stages 4-6
    # ═════════════════════════════════════════════════════════════════
    s = add_slide()
    add_bg(s, CREAM)
    section_header(s, "Methodology: Hydrology, AI & Drainage Design (Stages 4\u20136)")

    # Stage 4
    shape(s, Inches(0.6), Inches(1.7), Inches(3.8), Inches(2.7), fill=WHITE, line=NAVY)
    shape(s, Inches(0.6), Inches(1.7), Inches(3.8), Inches(0.45), fill=NAVY)
    _add(s, Inches(0.7), Inches(1.72), Inches(3.6), Inches(0.4),
         "Stage 4: Hydrological Analysis", size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    s4_items = [
        "Depression filling (Wang & Liu, pysheds)",
        "D8 flow direction \u2192 flow accumulation",
        "TWI = ln(\u03b1 / tan \u03b2)",
        "Stream extraction @ \u22651000 cell threshold",
        "O(N) DFS tracer for LineString extraction",
        "Output: twi.tif, flow_accum.tif, streams",
    ]
    add_bullets(s, Inches(0.7), Inches(2.3), Inches(3.6), Inches(2.0), s4_items, size=12, color=DARK, spacing=4)

    # Stage 5
    shape(s, Inches(4.65), Inches(1.7), Inches(3.8), Inches(2.7), fill=WHITE, line=NAVY)
    shape(s, Inches(4.65), Inches(1.7), Inches(3.8), Inches(0.45), fill=NAVY)
    _add(s, Inches(4.75), Inches(1.72), Inches(3.6), Inches(0.4),
         "Stage 5: Waterlogging Prediction", size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    s5_items = [
        "10-feature stack: elev, slope, aspect, TPI,",
        "  TWI, flow_acc, curvatures, depression, stream_dist",
        "Terrain-proxy labels: TWI>8 OR log_acc>7",
        "  OR low_elev+low_slope OR depression>0.2m",
        "XGBoost Classifier (100 trees, depth=6)",
        "5-fold StratifiedKFold CV + COG probability map",
    ]
    add_bullets(s, Inches(4.75), Inches(2.3), Inches(3.6), Inches(2.0), s5_items, size=12, color=DARK, spacing=4)

    # Stage 6
    shape(s, Inches(8.7), Inches(1.7), Inches(3.8), Inches(2.7), fill=WHITE, line=NAVY)
    shape(s, Inches(8.7), Inches(1.7), Inches(3.8), Inches(0.45), fill=NAVY)
    _add(s, Inches(8.8), Inches(1.72), Inches(3.6), Inches(0.4),
         "Stage 6: Drainage Design", size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    s6_items = [
        "Build flow graph from stream segments",
        "Minimum Spanning Tree (Kruskal\u2019s algorithm)",
        "Rational Method: Q = C\u00b7i\u00b7A / 360",
        "Manning\u2019s Eq: Q = (1/n)\u00b7A\u00b7R\u00b2\u00b3\u00b7S\u00b9\u00b2",
        "Trapezoidal earthen channel (1.5H:1V)",
        "Cost estimation per segment \u2192 GPKG export",
    ]
    add_bullets(s, Inches(8.8), Inches(2.3), Inches(3.6), Inches(2.0), s6_items, size=12, color=DARK, spacing=4)

    # Manning's equation highlight
    shape(s, Inches(0.6), Inches(4.7), Inches(11.9), Inches(2.3), fill=TAN, radius=0.05)
    _add(s, Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.4),
         "\U0001f4d0 Design Parameters (10-year return period)", size=16, bold=True, color=NAVY)
    design_params = [
        "Manning\u2019s n = 0.025 (earthen unlined)  \u00b7  Side slope z = 1.5 (1.5H:1V)  \u00b7  Rainfall i = 50 mm/hr  \u00b7  Runoff C = 0.65",
        "Velocity check: 0.3 m/s (self-cleaning) \u2264 V \u2264 2.0 m/s (erosion limit)  \u00b7  Cost rate: \u20b9800/m (base) + depth-dependent factor",
    ]
    add_bullets(s, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.3), design_params, size=13, color=DARK, symbol="\u25b8")

    # ═════════════════════════════════════════════════════════════════
    #  SLIDE 8 — ML MODEL ARCHITECTURE
    # ═════════════════════════════════════════════════════════════════
    s = add_slide()
    add_bg(s, CREAM)
    section_header(s, "ML Model: XGBoost Waterlogging Classifier")

    # Model architecture left
    shape(s, Inches(0.6), Inches(1.7), Inches(6.0), Inches(5.3), fill=WHITE, radius=0.05)
    _add(s, Inches(0.8), Inches(1.8), Inches(5.6), Inches(0.4),
         "Architecture & Training", size=18, bold=True, color=NAVY)
    model_items = [
        "Algorithm: XGBoost (Histogram-based tree splits)",
        "Estimators: 100 trees, max_depth=6",
        "Learning rate: 0.05, scale_pos_weight=5.5",
        "Cross-validation: 5-fold StratifiedKFold",
        "Preprocessing: RobustScaler (median/IQR)",
        "Training AUC: 1.000 (memorises terrain rule)",
        "CV AUC: 0.560 (DEVDI) / 0.509 (KHAPRETA)",
        "Model size: ~711 KB (joblib serialised)",
    ]
    add_bullets(s, Inches(0.8), Inches(2.3), Inches(5.6), Inches(4.5), model_items, size=13, color=DARK, spacing=5)

    # Feature importance right
    shape(s, Inches(6.9), Inches(1.7), Inches(5.8), Inches(5.3), fill=WHITE, radius=0.05)
    _add(s, Inches(7.1), Inches(1.8), Inches(5.4), Inches(0.4),
         "Feature Importance (DEVDI)", size=18, bold=True, color=NAVY)

    # Feature table
    feat_devdi = devdi.get("waterlogging", {}).get("feature_importances", [])
    feat_data = [["Feature", "Importance"]]
    for f in feat_devdi:
        feat_data.append([f["feature"], f"{f['importance']:.1%}"])
    make_table(s, Inches(7.1), Inches(2.4), Inches(5.4), Inches(3.5),
               len(feat_data), 2, feat_data,
               header_color=BROWN)

    # Key insight
    _add(s, Inches(7.1), Inches(6.0), Inches(5.4), Inches(0.8),
         "log_flow_accumulation dominates at 85\u201387% importance across both villages.\nTWI and stream_distance are secondary contributors.",
         size=12, color=SIENNA, italic=True)

    # ═════════════════════════════════════════════════════════════════
    #  SLIDE 9 — RESULTS DEVDI
    # ═════════════════════════════════════════════════════════════════
    s = add_slide()
    add_bg(s, CREAM)
    section_header(s, "Results: DEVDI Village", "64.6 M pts  \u00b7  0.81 km\u00b2  \u00b7  65 pts/m\u00b2  \u00b7  6 tiles  \u00b7  Elev: -31.2 to -7.9 m")

    # Left: Metrics table
    v = villages_data[0]
    gc = v.get("ground_classification", {})
    dt = v.get("dtm", {})
    wl = v.get("waterlogging", {})
    wm = wl.get("mean_metrics", {})
    dr = v.get("drainage", {})

    data = [
        ["Category", "Metric", "Value"],
        ["Ground Classification", "Accuracy / F1-Score", f"{gc.get('accuracy',0):.1%} / {gc.get('f1_score',0):.3f}"],
        ["", "Precision / Recall", f"{gc.get('precision',0):.3f} / {gc.get('recall',0):.1%}"],
        ["DTM Accuracy", "RMSE / MAE", f"{dt.get('rmse_m',0):.2f} m / {dt.get('mae_m',0):.2f} m"],
        ["", "LE90 / NMAD", f"{dt.get('le90_m',0):.2f} m / {dt.get('nmad_m',0):.3f} m"],
        ["Waterlogging", "Mean AUC (CV)", f"{wm.get('roc_auc',0):.4f}"],
        ["", "Mean F1 / Brier", f"{wm.get('f1',0):.4f} / {wl.get('brier_score',0):.4f}"],
        ["", "Positive Rate", f"{wl.get('positive_rate',0):.2%}"],
        ["Drainage Network", "Channels / Length", f"{dr.get('channel_count',0)} / {dr.get('total_length_m',0)/1000:.1f} km"],
        ["", "Total Cost", f"\u20b9{dr.get('total_cost_inr_lakhs',0):,.1f} Lakhs"],
        ["", "Avg Velocity", f"{dr.get('avg_velocity_ms',0):.3f} m/s"],
    ]
    make_table(s, Inches(0.5), Inches(1.8), Inches(7.5), Inches(5.2),
               len(data), 3, data,
               col_widths=[Inches(2.0), Inches(2.5), Inches(3.0)],
               header_color=NAVY)

    # Right: Image placeholder
    shape(s, Inches(8.3), Inches(1.8), Inches(4.5), Inches(3.5), fill=WHITE, line=GREY)
    shape(s, Inches(8.3), Inches(1.8), Inches(4.5), Inches(0.4), fill=NAVY)
    _add(s, Inches(8.4), Inches(1.82), Inches(4.3), Inches(0.35),
         "DTM + Drainage Map", size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    _add(s, Inches(9.0), Inches(3.2), Inches(3.0), Inches(0.5),
         "[Screenshot of DTM and\ndrainage network overlay]", size=13, color=GREY, align=PP_ALIGN.CENTER)

    # Waterlogging placeholder
    shape(s, Inches(8.3), Inches(5.5), Inches(4.5), Inches(1.8), fill=WHITE, line=GREY)
    shape(s, Inches(8.3), Inches(5.5), Inches(4.5), Inches(0.4), fill=NAVY)
    _add(s, Inches(8.4), Inches(5.52), Inches(4.3), Inches(0.35),
         "Waterlogging Probability Map", size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    _add(s, Inches(9.0), Inches(6.2), Inches(3.0), Inches(0.4),
         "[Screenshot of probability heatmap]", size=13, color=GREY, align=PP_ALIGN.CENTER)

    # ═════════════════════════════════════════════════════════════════
    #  SLIDE 10 — RESULTS KHAPRETA
    # ═════════════════════════════════════════════════════════════════
    s = add_slide()
    add_bg(s, CREAM)
    section_header(s, "Results: KHAPRETA Village", "163.7 M pts  \u00b7  0.67 km\u00b2  \u00b7  245 pts/m\u00b2  \u00b7  4 tiles  \u00b7  Elev: 115.3 to 151.0 m")

    v = villages_data[1]
    gc = v.get("ground_classification", {})
    dt = v.get("dtm", {})
    wl = v.get("waterlogging", {})
    wm = wl.get("mean_metrics", {})
    dr = v.get("drainage", {})

    data = [
        ["Category", "Metric", "Value"],
        ["Ground Classification", "Accuracy / F1-Score", f"{gc.get('accuracy',0):.1%} / {gc.get('f1_score',0):.3f}"],
        ["", "Precision / Recall", f"{gc.get('precision',0):.3f} / {gc.get('recall',0):.1%}"],
        ["DTM Accuracy", "RMSE / MAE", f"{dt.get('rmse_m',0):.2f} m / {dt.get('mae_m',0):.2f} m"],
        ["", "LE90 / NMAD", f"{dt.get('le90_m',0):.2f} m / {dt.get('nmad_m',0):.3f} m"],
        ["Waterlogging", "Mean AUC (CV)", f"{wm.get('roc_auc',0):.4f}"],
        ["", "Mean F1 / Brier", f"{wm.get('f1',0):.4f} / {wl.get('brier_score',0):.4f}"],
        ["", "Positive Rate", f"{wl.get('positive_rate',0):.2%}"],
        ["Drainage Network", "Channels / Length", f"{dr.get('channel_count',0)} / {dr.get('total_length_m',0)/1000:.1f} km"],
        ["", "Total Cost", f"\u20b9{dr.get('total_cost_inr_lakhs',0):,.1f} Lakhs"],
        ["", "Avg Velocity", f"{dr.get('avg_velocity_ms',0):.3f} m/s"],
    ]
    make_table(s, Inches(0.5), Inches(1.8), Inches(7.5), Inches(5.2),
               len(data), 3, data,
               col_widths=[Inches(2.0), Inches(2.5), Inches(3.0)],
               header_color=NAVY)

    shape(s, Inches(8.3), Inches(1.8), Inches(4.5), Inches(3.0), fill=WHITE, line=GREY)
    shape(s, Inches(8.3), Inches(1.8), Inches(4.5), Inches(0.4), fill=NAVY)
    _add(s, Inches(8.4), Inches(1.82), Inches(4.3), Inches(0.35),
         "DTM + Drainage Map", size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    _add(s, Inches(9.0), Inches(3.0), Inches(3.0), Inches(0.5),
         "[Screenshot]", size=13, color=GREY, align=PP_ALIGN.CENTER)

    shape(s, Inches(8.3), Inches(5.0), Inches(4.5), Inches(2.0), fill=WHITE, line=GREY)
    shape(s, Inches(8.3), Inches(5.0), Inches(4.5), Inches(0.4), fill=NAVY)
    _add(s, Inches(8.4), Inches(5.02), Inches(4.3), Inches(0.35),
         "Waterlogging Probability Map", size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    _add(s, Inches(9.0), Inches(5.8), Inches(3.0), Inches(0.4),
         "[Screenshot]", size=13, color=GREY, align=PP_ALIGN.CENTER)

    # ═════════════════════════════════════════════════════════════════
    #  SLIDE 11 — COMPARISON
    # ═════════════════════════════════════════════════════════════════
    s = add_slide()
    add_bg(s, CREAM)
    section_header(s, "Side-by-Side Village Comparison")

    data = [
        ["Metric", "DEVDI", "KHAPRETA"],
        ["Points", "64.6 M", "163.7 M"],
        ["Density", "65 pts/m\u00b2", "245 pts/m\u00b2"],
        ["Area", "0.81 km\u00b2", "0.67 km\u00b2"],
        ["Elevation Range", "-31.2 to -7.9 m", "115.3 to 151.0 m"],
        ["Relief", "23.3 m", "35.8 m"],
        ["GC F1-Score", f"{devdi['ground_classification']['f1_score']:.3f}", f"{khapreta['ground_classification']['f1_score']:.3f}"],
        ["DTM RMSE", f"{devdi['dtm']['rmse_m']:.2f} m", f"{khapreta['dtm']['rmse_m']:.2f} m"],
        ["DTM LE90", f"{devdi['dtm']['le90_m']:.2f} m", f"{khapreta['dtm']['le90_m']:.2f} m"],
        ["WL AUC", f"{devdi['waterlogging']['mean_metrics']['roc_auc']:.3f}", f"{khapreta['waterlogging']['mean_metrics']['roc_auc']:.3f}"],
        ["WL Positive Rate", f"{devdi['waterlogging']['positive_rate']:.1%}", f"{khapreta['waterlogging']['positive_rate']:.1%}"],
        ["Drain Channels", f"{devdi['drainage']['channel_count']}", f"{khapreta['drainage']['channel_count']}"],
        ["Drain Length", f"{devdi['drainage']['total_length_m']/1000:.1f} km", f"{khapreta['drainage']['total_length_m']/1000:.1f} km"],
        ["Total Cost", f"\u20b9{devdi['drainage']['total_cost_inr_lakhs']:,.1f} L", f"\u20b9{khapreta['drainage']['total_cost_inr_lakhs']:,.1f} L"],
    ]
    make_table(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.5),
               len(data), 3, data,
               col_widths=[Inches(3.0), Inches(4.65), Inches(4.65)],
               header_color=NAVY)

    # ═════════════════════════════════════════════════════════════════
    #  SLIDE 12 — COST ANALYSIS
    # ═════════════════════════════════════════════════════════════════
    s = add_slide()
    add_bg(s, CREAM)
    section_header(s, "Cost Analysis: Drainage Network Design")

    # Cost breakdown
    shape(s, Inches(0.6), Inches(1.7), Inches(6.0), Inches(5.2), fill=WHITE, radius=0.05)
    _add(s, Inches(0.8), Inches(1.8), Inches(5.6), Inches(0.4),
         "Cost Breakdown", size=18, bold=True, color=NAVY)

    # Calculate per-village costs
    total_cost = sum(v.get("drainage", {}).get("total_cost_inr_lakhs", 0) for v in villages_data)
    total_len = sum(v.get("drainage", {}).get("total_length_m", 0) for v in villages_data)
    total_ch = sum(v.get("drainage", {}).get("channel_count", 0) for v in villages_data)

    cost_items = [
        f"DEVDI: {devdi['drainage']['channel_count']} channels, {devdi['drainage']['total_length_m']/1000:.1f} km",
        f"  \u2192 Cost: \u20b9{devdi['drainage']['total_cost_inr_lakhs']:,.1f} Lakhs",
        f"  \u2192 Cost/km: \u20b9{devdi['drainage']['total_cost_inr_lakhs']/(devdi['drainage']['total_length_m']/1000):,.0f} L/km",
        "",
        f"KHAPRETA: {khapreta['drainage']['channel_count']} channels, {khapreta['drainage']['total_length_m']/1000:.1f} km",
        f"  \u2192 Cost: \u20b9{khapreta['drainage']['total_cost_inr_lakhs']:,.1f} Lakhs",
        f"  \u2192 Cost/km: \u20b9{khapreta['drainage']['total_cost_inr_lakhs']/(khapreta['drainage']['total_length_m']/1000):,.0f} L/km",
        "",
        f"TOTAL: {total_ch} channels, {total_len/1000:.1f} km across both villages",
        f"  \u2192 Combined Cost: \u20b9{total_cost:,.1f} Lakhs (\u20b9{total_cost/100:,.2f} Crores)",
    ]
    add_bullets(s, Inches(0.8), Inches(2.3), Inches(5.6), Inches(4.5), cost_items, size=14, color=DARK, spacing=4, symbol="")

    # Design summary right
    shape(s, Inches(6.9), Inches(1.7), Inches(5.8), Inches(5.2), fill=WHITE, radius=0.05)
    _add(s, Inches(7.1), Inches(1.8), Inches(5.4), Inches(0.4),
         "Design Parameters (both villages)", size=18, bold=True, color=NAVY)
    design_items = [
        "Design storm: 10-year return period",
        "Rainfall intensity: 50 mm/hr",
        "Runoff coefficient: C = 0.65 (Rational Method)",
        "Channel type: Trapezoidal earthen unlined",
        "Manning\u2019s n = 0.025",
        "Side slope: 1.5H : 1V",
        "Velocity limits: 0.3 \u2013 2.0 m/s",
        "", "Capacity exceedances: 0 (all channels adequate)",
        "Stream coverage: 100% (all extracted streams designed)",
    ]
    add_bullets(s, Inches(7.1), Inches(2.3), Inches(5.4), Inches(4.5), design_items, size=13, color=DARK, spacing=4)

    # ═════════════════════════════════════════════════════════════════
    #  SLIDE 13 — SCALABILITY
    # ═════════════════════════════════════════════════════════════════
    s = add_slide()
    add_bg(s, CREAM)
    section_header(s, "Scalability & Production Readiness")

    # Three columns
    cols = [
        ("\U0001f9ed Tiling Strategy", NAVY, [
            "Files >50M pts auto-tiled into",
            "500\u00d7500 m tiles with 25m buffer",
            "Each tile processed independently",
            "Parallelisable across CPU cores",
            "Merged after classification",
            "Memory-safe for any file size",
            "Tested: 64M \u2192 6 tiles, 163M \u2192 4 tiles",
        ]),
        ("\u26a1 Performance", BROWN, [
            "Full pipeline ~8 min per village",
            "Stage 1: 5s  |  Stage 2: 120s",
            "Stage 3: 25s  |  Stage 4: 50s",
            "Stage 5: 260s  |  Stage 6: 1s",
            "RAM: 8 GB min, 16 GB rec",
            "CPU: 4 cores min, 8+ rec",
            "Disk: 5 GB free per village",
        ]),
        ("\U0001f504 Batch Mode", SIENNA, [
            "--batch flag processes all villages",
            "Configured via config.yaml",
            "Rich console summary per village",
            "Per-village independent output dirs",
            "Can resume from any stage:",
            "  --stages 4,5,6 skips prior stages",
            "Designed for 10+ village rollout",
        ]),
    ]
    for i, (title, col_color, items) in enumerate(cols):
        x = Inches(0.5) + i * Inches(4.2)
        shape(s, x, Inches(1.7), Inches(3.9), Inches(5.3), fill=WHITE, line=col_color)
        shape(s, x, Inches(1.7), Inches(3.9), Inches(0.55), fill=col_color)
        _add(s, x + Inches(0.1), Inches(1.72), Inches(3.7), Inches(0.5),
             title, size=15, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_bullets(s, x + Inches(0.15), Inches(2.4), Inches(3.6), Inches(4.3),
                    items, size=12, color=DARK, spacing=5)

    # ═════════════════════════════════════════════════════════════════
    #  SLIDE 14 — KEY INSIGHTS
    # ═════════════════════════════════════════════════════════════════
    s = add_slide()
    add_bg(s, CREAM)
    section_header(s, "Key Insights & Learnings")

    insights = [
        "Both villages show ~15% waterlogging-prone area (consistent terrain-rule labelling across terrains)",
        "log_flow_accumulation is the dominant predictor (85\u201387% feature importance) \u2014 terrain convergence drives waterlogging",
        "KHAPRETA\u2019s higher relief (35.8 vs 23.3 m) drives steeper slopes, fewer channels (575 vs 966) and shorter drain length",
        "DEVDI at mean -18.3 m elevation suggests coastal/lowland terrain requiring denser drainage infrastructure",
        "Higher point density (245 pts/m\u00b2) captures more micro-terrain detail \u2192 higher DTM RMSE (6.13 vs 0.89 m) against flat-plane reference",
        "All drainage designs meet hydraulic capacity (0 capacity exceedances) with velocities within self-cleaning range",
        "Pipeline is fully automated: raw LAS/LAZ \u2192 engineered drainage network + GIS outputs in <10 minutes (64M points)",
    ]
    add_bullets(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.5), insights, size=15, color=DARK, spacing=7)

    # Next steps box
    shape(s, Inches(0.6), Inches(5.5), Inches(12.0), Inches(1.5), fill=TAN, radius=0.05)
    _add(s, Inches(0.8), Inches(5.6), Inches(11.5), Inches(0.4),
         "\U0001f3af Next Steps", size=16, bold=True, color=NAVY)
    next_items = [
        "Field-validate waterlogging hotspots  \u00b7  Obtain real flood labels (AUC expected >0.85)  \u00b7  Run batch on all 10 SVAMITVA villages",
        "DGPS checkpoint validation for DTM (target LE90 \u22640.15 m, ISRO standard)  \u00b7  Hydrograph routing for peak flow estimation",
    ]
    add_bullets(s, Inches(0.8), Inches(6.0), Inches(11.5), Inches(1.0), next_items, size=13, color=DARK, symbol="\u25b8")

    # ═════════════════════════════════════════════════════════════════
    #  SLIDE 15 — THANK YOU
    # ═════════════════════════════════════════════════════════════════
    s = add_slide()
    add_bg(s, CREAM)
    shape(s, 0, 0, W, Inches(0.08), fill=SIENNA)

    _add(s, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
         "Thank You", size=48, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    shape(s, Inches(5.5), Inches(3.1), Inches(2.5), Inches(0.05), fill=SIENNA)
    _add(s, Inches(1), Inches(3.3), Inches(11), Inches(0.6),
         "DTM Drainage AI Pipeline  \u2014  MoPR Geospatial Intelligence Hackathon",
         size=20, color=BROWN, align=PP_ALIGN.CENTER)
    _add(s, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
         "Problem Statement 2: DTM Generation, Waterlogging Prediction & Drainage Network Design",
         size=15, color=GREY, align=PP_ALIGN.CENTER)

    # Team
    shape(s, Inches(3.5), Inches(5.0), Inches(6.5), Inches(1.8), fill=NAVY, radius=0.06)
    _add(s, Inches(3.5), Inches(5.1), Inches(6.5), Inches(0.4),
         "Siddhartha Academy of Higher Education", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _add(s, Inches(3.5), Inches(5.55), Inches(6.5), Inches(1.1),
         "[Team Member 1]  \u00b7  [Team Member 2]  \u00b7  [Team Member 3]\n[Team Member 4]  \u00b7  [Team Member 5]",
         size=14, color=WHITE, align=PP_ALIGN.CENTER)

    # Save
    path = BASE / "DTM_Drainage_AI_Hackathon.pptx"
    prs.save(str(path))
    print(f"Presentation saved: {path}")
    return path

if __name__ == "__main__":
    make_pptx()
    print("Done!")
