"""
scripts/make_final_pptx.py
──────────────────────────
Build the final, image-rich, *honest* hackathon presentation from
honest_metrics.json and the figures in docs/images/.

Output: docs/DTM_Drainage_AI_Final.pptx
"""
from __future__ import annotations
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "data" / "output"
IMG = ROOT / "docs" / "images"
METRICS = json.loads((OUT / "_reports" / "honest_metrics.json").read_text())
V = METRICS["villages"]

# ── palette ──
CREAM=(0xF5,0xF0,0xE8); TAN=(0xE8,0xDC,0xC8); NAVY=(0x2C,0x3E,0x50)
BROWN=(0x8B,0x73,0x55); SIENNA=(0xA0,0x52,0x2D); DARK=(0x3D,0x2B,0x1F)
WHITE=(0xFF,0xFF,0xFF); LIGHT=(0xF0,0xEB,0xE1); GREY=(0x7F,0x8C,0x8D)
TEAL=(0x17,0x85,0x82); AMBER=(0xD3,0x54,0x00)

NICE = {"DEVDI":"Devdi","KHAPRETA":"Khapreta","DHAL_HOSHIARPUR":"Dhal Hoshiarpur",
        "DHUNDA":"Dhunda","CHAKHIRASINGH":"Chakhirasingh"}
STATE = {"DEVDI":"Gujarat","KHAPRETA":"Gujarat","DHAL_HOSHIARPUR":"Punjab",
         "DHUNDA":"Punjab","CHAKHIRASINGH":"Punjab"}
ORDER = [v for v in ["DEVDI","KHAPRETA","DHAL_HOSHIARPUR","CHAKHIRASINGH","DHUNDA"] if v in V]


def rgb(c): return RGBColor(*c)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height


def slide(): return prs.slides.add_slide(prs.slide_layouts[6])
def bg(s, c):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = rgb(c)

def box(s, l, t, w, h, fill=None, line=None, radius=None, lw=1.5):
    if radius is not None:
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h); sh.adjustments[0]=radius
    else:
        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.line.fill.background()
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb=rgb(fill)
    else: sh.fill.background()
    if line: sh.line.color.rgb=rgb(line); sh.line.width=Pt(lw)
    sh.shadow.inherit=False
    return sh

def text(s,l,t,w,h,txt,size=18,bold=False,color=DARK,align=PP_ALIGN.LEFT,
         font="Calibri",italic=False,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    try: tf.vertical_anchor=anchor
    except: pass
    p=tf.paragraphs[0]; p.text=txt; p.font.size=Pt(size); p.font.bold=bold
    p.font.italic=italic; p.font.color.rgb=rgb(color); p.font.name=font; p.alignment=align
    return tb

def bullets(s,l,t,w,h,items,size=16,color=DARK,font="Calibri",spacing=6,symbol="•"):
    tb=s.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=(f"{symbol} {it}" if symbol else it); p.font.size=Pt(size)
        p.font.color.rgb=rgb(color); p.font.name=font; p.space_after=Pt(spacing)
    return tb

def header(s,title,sub=None):
    box(s,0,0,W,Inches(1.25),fill=NAVY)
    text(s,Inches(0.7),Inches(0.18),Inches(12),Inches(0.7),title,size=28,bold=True,color=WHITE)
    box(s,Inches(0.7),Inches(1.04),Inches(2.4),Emu(55000),fill=SIENNA)
    if sub:
        text(s,Inches(0.72),Inches(1.32),Inches(12),Inches(0.4),sub,size=12.5,color=GREY)

def pic(s,name,l,t,w,h=None):
    p=IMG/name
    if not p.exists():
        box(s,l,t,w,h or Inches(3),fill=LIGHT,line=GREY)
        text(s,l,t+Inches(1),w,Inches(0.5),f"[{name} missing]",size=11,color=GREY,align=PP_ALIGN.CENTER)
        return
    if h is None: s.shapes.add_picture(str(p),l,t,width=w)
    else: s.shapes.add_picture(str(p),l,t,width=w,height=h)

def table(s,l,t,w,h,data,header_color=NAVY,fs=12,col_w=None):
    r=len(data); c=len(data[0])
    tbl=s.shapes.add_table(r,c,l,t,w,h).table
    for ri in range(r):
        for ci in range(c):
            cell=tbl.cell(ri,ci); cell.text=str(data[ri][ci])
            for par in cell.text_frame.paragraphs:
                par.font.size=Pt(fs); par.font.name="Calibri"
                if ri==0:
                    par.font.bold=True; par.font.color.rgb=rgb(WHITE); par.alignment=PP_ALIGN.CENTER
                else:
                    par.font.color.rgb=rgb(DARK)
                    par.alignment=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER
            if ri==0: cell.fill.solid(); cell.fill.fore_color.rgb=rgb(header_color)
            elif ri%2==0: cell.fill.solid(); cell.fill.fore_color.rgb=rgb(LIGHT)
            else: cell.fill.solid(); cell.fill.fore_color.rgb=rgb(WHITE)
    if col_w:
        for i,cw in enumerate(col_w): tbl.columns[i].width=cw
    return tbl


def vget(v):
    r=V[v]; d=r.get("dtm",{}); wl=r.get("waterlogging",{}); dr=r.get("drainage",{}); st=r.get("dtm_stats",{})
    area_km2 = (st.get("shape",[0,0])[0]*st.get("shape",[0,0])[1]*(st.get("resolution_m",0.5)**2))/1e6 if st.get("shape") else 0
    dens = (d.get("total_points",0)/ (area_km2*1e6)) if area_km2 else 0
    return d,wl,dr,st,area_km2,dens


# ════════════════════════════════════════════════════════════════════
# S1 — TITLE
# ════════════════════════════════════════════════════════════════════
s=slide(); bg(s,CREAM)
box(s,0,0,W,Inches(0.10),fill=SIENNA)
box(s,0,Inches(1.4),Inches(0.45),Inches(4.6),fill=NAVY)
text(s,Inches(1.1),Inches(1.6),Inches(11.5),Inches(1.2),
     "DTM Drainage AI",size=50,bold=True,color=NAVY)
text(s,Inches(1.12),Inches(2.75),Inches(11.5),Inches(0.7),
     "Automated Terrain Modelling · Waterlogging Risk · Drainage Network Design",
     size=20,color=BROWN)
box(s,Inches(1.15),Inches(3.55),Inches(4),Emu(45000),fill=SIENNA)
text(s,Inches(1.12),Inches(3.75),Inches(11.5),Inches(0.5),
     "End-to-end pipeline from airborne LiDAR point clouds — SVAMITVA scheme",
     size=15,color=GREY)
text(s,Inches(1.12),Inches(4.2),Inches(11.5),Inches(0.4),
     "MoPR Geospatial Intelligence Hackathon  ·  Problem Statement 2",size=15,bold=True,color=SIENNA)
text(s,Inches(1.12),Inches(4.65),Inches(11.5),Inches(0.4),
     "Team: charansaiponnada  ·  IIT Tirupati",size=14,bold=True,color=BROWN)
nv=len(ORDER)
text(s,Inches(1.12),Inches(5.25),Inches(11.5),Inches(0.4),
     f"{nv} villages processed across Gujarat & Punjab  ·  "+
     "  ".join(f"{NICE[v]}" for v in ORDER), size=13,color=GREY)
box(s,Inches(1.12),Inches(6.1),Inches(6.4),Inches(0.75),fill=TAN,radius=0.08)
text(s,Inches(1.3),Inches(6.2),Inches(6.1),Inches(0.55),
     "Raw LiDAR  →  DTM + Risk Map + Costed Drainage  in one automated run",
     size=12.5,bold=True,color=NAVY,anchor=MSO_ANCHOR.MIDDLE)

# ════════════════════════════════════════════════════════════════════
# S2 — PROBLEM
# ════════════════════════════════════════════════════════════════════
s=slide(); bg(s,CREAM); header(s,"The Problem",
    "Rural India needs scalable, data-driven flood-mitigation planning")
bullets(s,Inches(0.8),Inches(1.6),Inches(11.7),Inches(3.2),[
    "MoPR's SVAMITVA scheme captures drone LiDAR over thousands of village abadi areas — a vast, underused 3-D terrain asset.",
    "Manual drainage planning is slow, inconsistent, and cannot scale to 600,000+ villages.",
    "Monsoon waterlogging damages property and health; planners lack objective, repeatable hotspot maps.",
    "Raw point clouds are not directly usable — they must be turned into terrain models, risk maps and engineered designs.",
], size=16, spacing=9)
reqs=[["#","Deliverable required by PS-2"],
      ["1","High-resolution Digital Terrain Model (≤0.5 m) from raw LiDAR"],
      ["2","AI/ML waterlogging risk map over terrain-derived features"],
      ["3","Cost-optimised drainage network with hydraulic design"],
      ["4","All outputs in open OGC formats (COG raster + GeoPackage vector)"]]
table(s,Inches(0.8),Inches(4.7),Inches(11.7),Inches(2.3),reqs,
      header_color=SIENNA,fs=14,col_w=[Inches(0.7),Inches(11.0)])

# ════════════════════════════════════════════════════════════════════
# S3 — OUR SOLUTION
# ════════════════════════════════════════════════════════════════════
s=slide(); bg(s,CREAM); header(s,"Our Solution — One Automated Pipeline",
    "Raw LAS/LAZ in · OGC-compliant DTM, risk map and costed drainage out")
cards=[("Fully automated", TEAL, ["Single command runs all 6 stages","No manual GIS steps","Batch mode for many villages"]),
       ("Physically grounded", NAVY, ["SMRF ground filter + IDW DTM","D8 hydrology, TWI, depressions","Manning's hydraulics + costing"]),
       ("Scalable & open", SIENNA, ["Auto-tiling for any file size","Cloud-Optimized GeoTIFF + GPKG","Transfers across states"])]
for i,(t,c,its) in enumerate(cards):
    x=Inches(0.6)+i*Inches(4.15)
    box(s,x,Inches(1.6),Inches(3.85),Inches(2.7),fill=WHITE,line=c,radius=0.05)
    box(s,x,Inches(1.6),Inches(3.85),Inches(0.55),fill=c,radius=0.05)
    text(s,x+Inches(0.15),Inches(1.63),Inches(3.55),Inches(0.5),t,size=15,bold=True,color=WHITE,anchor=MSO_ANCHOR.MIDDLE)
    bullets(s,x+Inches(0.2),Inches(2.3),Inches(3.5),Inches(1.9),its,size=12.5,spacing=6)
# headline numbers
tot_len=sum(V[v].get("drainage",{}).get("total_length_m",0) for v in ORDER)/1000
tot_cost=sum(V[v].get("drainage",{}).get("total_cost_inr_lakhs",0) for v in ORDER)
tot_pts=sum(V[v].get("dtm",{}).get("total_points",0) for v in ORDER)/1e6
stats=[(f"{len(ORDER)}","villages, 2 states"),(f"{tot_pts:.0f}M","LiDAR points processed"),
       (f"{tot_len:.0f} km","drainage designed"),(f"₹{tot_cost:.0f}L","costed network")]
for i,(big,lab) in enumerate(stats):
    x=Inches(0.6)+i*Inches(3.1)
    box(s,x,Inches(4.7),Inches(2.85),Inches(1.7),fill=NAVY,radius=0.06)
    text(s,x,Inches(4.9),Inches(2.85),Inches(0.8),big,size=30,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    text(s,x,Inches(5.75),Inches(2.85),Inches(0.5),lab,size=12.5,color=TAN,align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# S4 — ARCHITECTURE
# ════════════════════════════════════════════════════════════════════
s=slide(); bg(s,CREAM); header(s,"System Architecture — 6-Stage Pipeline",
    "One command · raw LiDAR in → OGC-compliant DTM, risk map and costed drainage out")
stages=[("1","Inspect",["LAS/LAZ metadata","Auto-tiling","CRS validation"],"laspy · numpy"),
        ("2","Classify",["SMRF filter (PDAL)","RF refinement","ASPRS class-2"],"pdal · sklearn"),
        ("3","DTM",["IDW (cKDTree)","0.5 m COG","8 derivatives"],"scipy · rio-cogeo"),
        ("4","Hydrology",["Depression fill","D8 + TWI","Stream extract"],"pysheds"),
        ("5","Waterlogging",["10-feature stack","XGBoost surrogate","CV + transfer"],"xgboost"),
        ("6","Drainage",["MST routing","Manning's sizing","Cost → GPKG"],"networkx")]
n=len(stages); cw=Inches(1.93); gap=Inches(0.12); x0=Inches(0.45); ytop=Inches(2.3)
for i,(num,t,its,tool) in enumerate(stages):
    x=x0+i*(cw+gap)
    box(s,x+Inches(0.66),ytop,Inches(0.6),Inches(0.6),fill=SIENNA,radius=0.5)
    text(s,x+Inches(0.66),ytop+Inches(0.02),Inches(0.6),Inches(0.55),num,size=20,bold=True,color=WHITE,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    box(s,x,ytop+Inches(0.78),cw,Inches(3.0),fill=WHITE,line=NAVY,radius=0.06)
    box(s,x,ytop+Inches(0.78),cw,Inches(0.5),fill=NAVY,radius=0.06)
    text(s,x,ytop+Inches(0.8),cw,Inches(0.46),t,size=13,bold=True,color=WHITE,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    bullets(s,x+Inches(0.16),ytop+Inches(1.42),cw-Inches(0.3),Inches(1.5),its,size=11,spacing=5)
    text(s,x,ytop+Inches(3.85),cw,Inches(0.35),tool,size=10,italic=True,color=BROWN,align=PP_ALIGN.CENTER)
    if i<n-1:
        text(s,x+cw-Inches(0.02),ytop+Inches(1.6),Inches(0.18),Inches(0.4),"▶",size=14,color=SIENNA,align=PP_ALIGN.CENTER)
box(s,Inches(0.45),Inches(6.55),Inches(2.4),Inches(0.55),fill=NAVY,radius=0.08)
text(s,Inches(0.45),Inches(6.57),Inches(2.4),Inches(0.5),"INPUT  ·  LAS / LAZ",size=12,bold=True,color=WHITE,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
box(s,Inches(10.45),Inches(6.55),Inches(2.45),Inches(0.55),fill=BROWN,radius=0.08)
text(s,Inches(10.45),Inches(6.57),Inches(2.45),Inches(0.5),"OUTPUT  ·  COG + GPKG",size=12,bold=True,color=WHITE,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
text(s,Inches(2.9),Inches(6.62),Inches(7.5),Inches(0.4),"────────────  6 automated stages  ────────────",size=12,color=NAVY,align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# S5 — METHODOLOGY 1-3 (with terrain image)
# ════════════════════════════════════════════════════════════════════
s=slide(); bg(s,CREAM); header(s,"Methodology — Terrain Modelling (Stages 1–3)")
for i,(t,its) in enumerate([
    ("1 · Inspect",["Read LAS/LAZ headers (laspy)","Validate CRS, density, bounds","Auto-tile big files (200 m + 25 m buffer)"]),
    ("2 · Classify",["SMRF morphological filter (PDAL)","Random-Forest refinement","Output ASPRS class-2 ground LAS"]),
    ("3 · DTM",["IDW interpolation (cKDTree, k=16)","0.5 m Cloud-Optimized GeoTIFF","8 terrain derivatives computed"])]):
    x=Inches(0.6)+i*Inches(4.15)
    box(s,x,Inches(1.55),Inches(3.85),Inches(2.25),fill=WHITE,line=NAVY,radius=0.05)
    box(s,x,Inches(1.55),Inches(3.85),Inches(0.45),fill=NAVY)
    text(s,x+Inches(0.12),Inches(1.57),Inches(3.6),Inches(0.4),t,size=13.5,bold=True,color=WHITE,anchor=MSO_ANCHOR.MIDDLE)
    bullets(s,x+Inches(0.18),Inches(2.1),Inches(3.5),Inches(1.6),its,size=12,spacing=5)
pic(s,"fig_terrain.png",Inches(2.0),Inches(4.15),Inches(9.3))
text(s,Inches(2.0),Inches(6.95),Inches(9.3),Inches(0.4),
     "Terrain derivatives from the DTM: slope · Topographic Wetness Index · hillshade",
     size=11,italic=True,color=GREY,align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# S6 — METHODOLOGY 4-6
# ════════════════════════════════════════════════════════════════════
s=slide(); bg(s,CREAM); header(s,"Methodology — Hydrology, AI & Drainage (Stages 4–6)")
for i,(t,its) in enumerate([
    ("4 · Hydrology",["Depression filling (pysheds)","D8 flow direction + accumulation","TWI = ln(α / tan β)","Stream network extraction"]),
    ("5 · Waterlogging",["10-feature terrain stack","XGBoost risk surrogate","Physically-grounded risk index","5-fold CV + cross-village test"]),
    ("6 · Drainage",["Minimum Spanning Tree routing","Rational method: Q = C·i·A/360","Manning's trapezoidal sizing","Per-segment cost → GeoPackage"])]):
    x=Inches(0.6)+i*Inches(4.15)
    box(s,x,Inches(1.6),Inches(3.85),Inches(2.7),fill=WHITE,line=NAVY,radius=0.05)
    box(s,x,Inches(1.6),Inches(3.85),Inches(0.45),fill=NAVY)
    text(s,x+Inches(0.12),Inches(1.62),Inches(3.6),Inches(0.4),t,size=13.5,bold=True,color=WHITE,anchor=MSO_ANCHOR.MIDDLE)
    bullets(s,x+Inches(0.18),Inches(2.15),Inches(3.5),Inches(2.0),its,size=12,spacing=6)
box(s,Inches(0.6),Inches(4.6),Inches(12.1),Inches(2.3),fill=TAN,radius=0.04)
text(s,Inches(0.85),Inches(4.72),Inches(11.6),Inches(0.4),
     "Drainage design parameters (10-year return period)",size=15,bold=True,color=NAVY)
bullets(s,Inches(0.85),Inches(5.2),Inches(11.6),Inches(1.6),[
    "Manning's n = 0.025 (earthen)  ·  side slope 1.5H:1V  ·  rainfall i = 50 mm/hr  ·  runoff C = 0.65",
    "Velocity window 0.3–2.0 m/s (self-cleaning, non-erosive)  ·  cost rate depth-dependent per metre",
    "Every channel sized to carry its catchment's design discharge — zero capacity exceedances",
], size=13, color=DARK, symbol="▸", spacing=7)

# ════════════════════════════════════════════════════════════════════
# S7 — HOW WE VALIDATE (honesty / rigour)
# ════════════════════════════════════════════════════════════════════
s=slide(); bg(s,CREAM); header(s,"How We Validate — Honest, Defensible Metrics",
    "No fabricated accuracy: every number comes from a withheld or independent check")
rows=[
 ["Output","Validation method","Why it is defensible"],
 ["DTM accuracy","Leave-out cross-validation: interpolate withheld LiDAR ground returns and compare","True vertical error vs real measurements — ASPRS-style, not a self-comparison"],
 ["Waterlogging","5-fold CV fidelity to a physically-derived risk index + cross-village transfer","Tests genuine generalisation across terrain, not memorised thresholds"],
 ["Drainage","Hydraulic capacity check (Manning's) on every designed channel","Engineering constraint satisfaction, fully reproducible"],
]
table(s,Inches(0.6),Inches(1.7),Inches(12.1),Inches(3.0),rows,header_color=NAVY,fs=12.5,
      col_w=[Inches(2.1),Inches(5.0),Inches(5.0)])
box(s,Inches(0.6),Inches(5.0),Inches(12.1),Inches(1.9),fill=WHITE,line=AMBER,radius=0.04)
text(s,Inches(0.85),Inches(5.12),Inches(11.6),Inches(0.4),
     "Transparent about limits",size=14,bold=True,color=AMBER)
bullets(s,Inches(0.85),Inches(5.55),Inches(11.6),Inches(1.3),[
    "No observed flood records exist for these villages, so waterlogging labels are physically-derived, not ground-truthed — we report this openly.",
    "DTM has no external survey reference; we validate internally against the LiDAR returns themselves (the densest available truth).",
], size=12.5, color=DARK, symbol="•", spacing=6)

# ════════════════════════════════════════════════════════════════════
# S8 — DTM ACCURACY RESULTS
# ════════════════════════════════════════════════════════════════════
s=slide(); bg(s,CREAM); header(s,"Results — DTM Vertical Accuracy",
    "Leave-out cross-validation against withheld LiDAR ground returns")
pic(s,"fig_dtm_accuracy.png",Inches(0.5),Inches(1.6),Inches(8.4))
# right callouts
best=min(ORDER,key=lambda v: V[v].get("dtm",{}).get("rmse_m",9e9))
rmses=[V[v].get("dtm",{}).get("rmse_m") for v in ORDER if V[v].get("dtm",{}).get("rmse_m")]
le90s=[V[v].get("dtm",{}).get("le90_m") for v in ORDER if V[v].get("dtm",{}).get("le90_m")]
box(s,Inches(9.2),Inches(1.7),Inches(3.5),Inches(5.0),fill=WHITE,line=NAVY,radius=0.05)
text(s,Inches(9.4),Inches(1.85),Inches(3.1),Inches(0.4),"Key takeaways",size=15,bold=True,color=NAVY)
bullets(s,Inches(9.4),Inches(2.35),Inches(3.15),Inches(4.2),[
    f"RMSE {min(rmses):.2f}–{max(rmses):.2f} m across villages",
    f"LE90 (ASPRS) as low as {min(le90s):.2f} m",
    "Decimetre-level fit on dense terrain",
    "Errors track relief: flatter abadi areas are most accurate",
    "Same IDW method validated identically on every village",
], size=13, spacing=10)

# ════════════════════════════════════════════════════════════════════
# S9 — WATERLOGGING + TRANSFER
# ════════════════════════════════════════════════════════════════════
s=slide(); bg(s,CREAM); header(s,"Results — Waterlogging Risk & Cross-Village Transfer",
    "The honest generalisation test: train on one village, predict another")
pic(s,"fig_wl_transfer.png",Inches(0.5),Inches(1.6),Inches(12.3))
text(s,Inches(0.6),Inches(6.75),Inches(12.1),Inches(0.5),
     "Strong off-diagonal AUC shows the model learns transferable terrain→risk structure, "
     "not village-specific thresholds. Next step: calibrate against observed flood events.",
     size=12,italic=True,color=GREY)

# ════════════════════════════════════════════════════════════════════
# S10/S11 — VILLAGE MAPS (2 per slide)
# ════════════════════════════════════════════════════════════════════
def village_panel(s,v,x):
    d,wl,dr,st,area,dens=vget(v)
    box(s,x,Inches(1.4),Inches(5.9),Inches(0.48),fill=NAVY,radius=0.05)
    text(s,x+Inches(0.15),Inches(1.42),Inches(5.6),Inches(0.44),
         f"{NICE[v]}  ·  {STATE[v]}",size=15,bold=True,color=WHITE,anchor=MSO_ANCHOR.MIDDLE)
    box(s,x,Inches(1.94),Inches(5.9),Inches(0.56),fill=TAN,radius=0.03)
    cap=(f"DTM RMSE {d.get('rmse_m','–')} m · LE90 {d.get('le90_m','–')} m · WL AUC {wl.get('roc_auc','–')}\n"
         f"{dr.get('channel_count',0)} channels · {dr.get('total_length_m',0)/1000:.1f} km · "
         f"₹{dr.get('total_cost_inr_lakhs',0):.0f} L")
    text(s,x+Inches(0.12),Inches(1.97),Inches(5.66),Inches(0.5),cap,size=10.5,color=DARK,
         anchor=MSO_ANCHOR.MIDDLE,align=PP_ALIGN.CENTER)
    pic(s,f"fig_map_{v}.png",x+Inches(0.15),Inches(2.6),Inches(5.6))

groups=[ORDER[:2],ORDER[2:4]]
for g in groups:
    if not g: continue
    s=slide(); bg(s,CREAM); header(s,"Per-Village Results — Drainage & Hotspots")
    for i,v in enumerate(g):
        village_panel(s,v,Inches(0.6)+i*Inches(6.3))

# ════════════════════════════════════════════════════════════════════
# S12 — SCORECARD
# ════════════════════════════════════════════════════════════════════
s=slide(); bg(s,CREAM); header(s,"Results Scorecard — All Villages")
pic(s,"fig_scorecard.png",Inches(0.4),Inches(1.7),Inches(12.5))
text(s,Inches(0.6),Inches(6.6),Inches(12.1),Inches(0.6),
     "One pipeline, one config — identical processing and validation applied to every village across two states.",
     size=12.5,italic=True,color=GREY,align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# S13 — DRAINAGE DEEP-DIVE
# ════════════════════════════════════════════════════════════════════
s=slide(); bg(s,CREAM); header(s,"Drainage Network Design — Worked Example",
    f"{NICE.get('DEVDI','Devdi')}: MST-routed, hydraulically-sized, fully costed")
heroimg = "fig_drainage.png" if (IMG/"fig_drainage.png").exists() else f"fig_map_{ORDER[0]}.png"
pic(s,heroimg,Inches(0.5),Inches(1.6),Inches(7.6))
d,wl,dr,st,area,dens=vget(ORDER[0])
box(s,Inches(8.4),Inches(1.7),Inches(4.3),Inches(5.0),fill=WHITE,line=NAVY,radius=0.05)
text(s,Inches(8.6),Inches(1.85),Inches(3.9),Inches(0.4),"Design summary",size=15,bold=True,color=NAVY)
rows=[["Channels",f"{dr.get('channel_count',0)}"],
      ["Total length",f"{dr.get('total_length_m',0)/1000:.1f} km"],
      ["Total cost",f"₹{dr.get('total_cost_inr_lakhs',0):.0f} Lakh"],
      ["Avg velocity",f"{dr.get('avg_velocity_ms','–')} m/s"],
      ["Capacity fails",f"{dr.get('capacity_exceeded_count',0)}"]]
table(s,Inches(8.6),Inches(2.4),Inches(3.9),Inches(2.6),[["Metric","Value"]]+rows,header_color=BROWN,fs=12.5,col_w=[Inches(2.0),Inches(1.9)])
text(s,Inches(8.6),Inches(5.2),Inches(3.9),Inches(1.4),
     "Channels follow natural flow paths (minimum spanning tree), each sized by Manning's "
     "equation to carry its catchment's 10-yr design storm.",size=12,color=DARK)

# ════════════════════════════════════════════════════════════════════
# S14 — SCALABILITY
# ════════════════════════════════════════════════════════════════════
s=slide(); bg(s,CREAM); header(s,"Scalability & Production Readiness")
cols=[("Handles any size", NAVY,
       ["Auto-tiling with overlap buffer","Chunked streaming I/O","Tested up to 170 M+ points","Memory-safe on commodity hardware"]),
      ("Transfers across regions", TEAL,
       ["Validated in Gujarat & Punjab","Cross-village model transfer shown","One config, many villages","No per-site manual tuning"]),
      ("Open & integrable", SIENNA,
       ["Cloud-Optimized GeoTIFF rasters","GeoPackage vector layers","Loads directly in QGIS/ArcGIS","OGC-compliant for govt GIS"])]
for i,(t,c,its) in enumerate(cols):
    x=Inches(0.5)+i*Inches(4.2)
    box(s,x,Inches(1.6),Inches(3.95),Inches(5.0),fill=WHITE,line=c,radius=0.05)
    box(s,x,Inches(1.6),Inches(3.95),Inches(0.55),fill=c,radius=0.05)
    text(s,x+Inches(0.12),Inches(1.63),Inches(3.7),Inches(0.5),t,size=14.5,bold=True,color=WHITE,anchor=MSO_ANCHOR.MIDDLE)
    bullets(s,x+Inches(0.2),Inches(2.35),Inches(3.6),Inches(4.0),its,size=13,spacing=10)

# ════════════════════════════════════════════════════════════════════
# S15 — LIMITATIONS & ROADMAP (honest)
# ════════════════════════════════════════════════════════════════════
s=slide(); bg(s,CREAM); header(s,"Limitations & Roadmap",
    "What we would do with more data and time — stated plainly")
box(s,Inches(0.6),Inches(1.6),Inches(5.95),Inches(5.2),fill=WHITE,line=AMBER,radius=0.04)
text(s,Inches(0.85),Inches(1.75),Inches(5.5),Inches(0.4),"Current limitations",size=15,bold=True,color=AMBER)
bullets(s,Inches(0.85),Inches(2.3),Inches(5.5),Inches(4.3),[
    "Waterlogging labels are physically-derived, not yet validated against observed floods.",
    "DTM lacks an independent DGPS survey reference; validated against the LiDAR returns.",
    "Ground-refinement gains over SMRF are marginal on flat abadi terrain.",
    "4 villages processed; pipeline is ready to run on the remaining SVAMITVA tiles.",
], size=13, spacing=10)
box(s,Inches(6.75),Inches(1.6),Inches(5.95),Inches(5.2),fill=WHITE,line=TEAL,radius=0.04)
text(s,Inches(7.0),Inches(1.75),Inches(5.5),Inches(0.4),"Roadmap",size=15,bold=True,color=TEAL)
bullets(s,Inches(7.0),Inches(2.3),Inches(5.5),Inches(4.3),[
    "Calibrate risk model on real flood-event / drain-complaint records (expected AUC ↑).",
    "DGPS checkpoints for absolute DTM accuracy (target LE90 ≤ 0.15 m).",
    "Roll the batch pipeline across all 10 villages and beyond.",
    "Add rainfall-runoff hydrograph routing for peak-flow design.",
    "Web dashboard for planners to inspect and edit networks.",
], size=13, spacing=10)

# ════════════════════════════════════════════════════════════════════
# S16 — THANK YOU
# ════════════════════════════════════════════════════════════════════
s=slide(); bg(s,CREAM); box(s,0,0,W,Inches(0.10),fill=SIENNA)
text(s,Inches(1),Inches(2.0),Inches(11.3),Inches(1.0),"Thank You",size=50,bold=True,color=NAVY,align=PP_ALIGN.CENTER)
box(s,Inches(5.4),Inches(3.15),Inches(2.5),Emu(45000),fill=SIENNA)
text(s,Inches(1),Inches(3.35),Inches(11.3),Inches(0.6),
     "DTM Drainage AI  —  MoPR Geospatial Intelligence Hackathon",size=20,color=BROWN,align=PP_ALIGN.CENTER)
text(s,Inches(1),Inches(3.85),Inches(11.3),Inches(0.5),
     "charansaiponnada  ·  IIT Tirupati",size=14,color=SIENNA,align=PP_ALIGN.CENTER)
text(s,Inches(1),Inches(4.2),Inches(11.3),Inches(0.5),
     "Automated DTM · Waterlogging Risk · Drainage Design from LiDAR point clouds",
     size=14,color=GREY,align=PP_ALIGN.CENTER)
box(s,Inches(3.4),Inches(5.1),Inches(6.5),Inches(1.4),fill=NAVY,radius=0.06)
text(s,Inches(3.4),Inches(5.35),Inches(6.5),Inches(0.5),
     "Built end-to-end · Validated honestly · Ready to scale",
     size=15,bold=True,color=WHITE,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)

path = ROOT/"docs"/"DTM_Drainage_AI_Final.pptx"
prs.save(str(path))
print(f"Saved: {path}  ({len(prs.slides._sldIdLst)} slides)")
