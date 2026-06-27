"""
scripts/pptx_results.py
───────────────────────
Post-process the built deck (DTM_Hackathon_Final.pptx) to spread the results
across multiple slides WITH real result-map images, instead of cramming
everything onto one slide.

  • Slide 5  "Results & Demonstration"  → de-crammed to DTM accuracy + highlights
  • NEW 6    "Waterlogging Risk Maps"   → risk_grid.png + model scores + transfer
  • NEW 7    "Drainage Network Design"  → result map + 4-village cost table

Works directly on the .pptx (the original inject/build scripts were temp files).
Run:  .venv/Scripts/python.exe scripts/pptx_results.py
"""
from __future__ import annotations
import copy, json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

ROOT = Path(__file__).resolve().parents[1]
DECK = ROOT / "DTM_Hackathon_Final.pptx"
FIG  = ROOT / "data" / "output" / "figures"
MET  = json.loads((ROOT / "data" / "output" / "_reports" / "honest_metrics.json").read_text())

RED=RGBColor(0x85,0x20,0x0C); NAVY=RGBColor(0x1B,0x3A,0x5C); TEAL=RGBColor(0x11,0x77,0x77)
AMBER=RGBColor(0xB9,0x4A,0x00); WHITE=RGBColor(0xFF,0xFF,0xFF); DARK=RGBColor(0x1A,0x1A,0x2E)
GREY=RGBColor(0x5D,0x6D,0x7E); LT=RGBColor(0xEA,0xF0,0xF6)

VK = ["DEVDI","KHAPRETA","DHAL_HOSHIARPUR","CHAKHIRASINGH"]
VLAB = {"DEVDI":"Devdi","KHAPRETA":"Khapreta","DHAL_HOSHIARPUR":"Dhal Hsp.","CHAKHIRASINGH":"Chakhira"}


def box(slide, x, y, w, h, paras, fill=None, line=None, anchor=MSO_ANCHOR.TOP):
    """paras = list of (text, size_pt, bold, color, align). Returns the shape."""
    sp = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = sp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(7)
    tf.margin_top = tf.margin_bottom = Pt(4)
    for i,(text,sz,bold,color,align) in enumerate(paras):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(sz); r.font.bold = bold; r.font.color.rgb = color
        r.font.name = "Calibri"
    if fill is not None:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    else:
        sp.fill.background()
    if line is not None:
        sp.line.color.rgb = line; sp.line.width = Pt(1.2)
    else:
        sp.line.fill.background()
    return sp


def pic_fit(slide, path, x, y, max_w, max_h):
    """Place image fit within (max_w,max_h) keeping aspect, centred in the box."""
    iw, ih = Image.open(path).size
    ar = iw/ih
    w, h = max_w, max_w/ar
    if h > max_h:
        h, w = max_h, max_h*ar
    px = x + (max_w - w)/2
    py = y + (max_h - h)/2
    slide.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(w), Inches(h))
    return px, py, w, h


def dup_results_slide(prs, src_idx, new_title):
    """Duplicate a slide keeping only its placeholders (red header + footer),
    retitle, and return the fresh slide (appended at end)."""
    src = prs.slides[src_idx]
    s = prs.slides.add_slide(src.slide_layout)
    for sh in list(s.shapes):                       # clear layout placeholders
        sh._element.getparent().remove(sh._element)
    for sh in src.shapes:                           # copy header/footer placeholders
        if sh.is_placeholder:
            s.shapes._spTree.append(copy.deepcopy(sh._element))
    # retitle the header placeholder (the one that isn't the footer / slide number)
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip().startswith("Results"):
            runs = [r for p in sh.text_frame.paragraphs for r in p.runs]
            if runs:
                runs[0].text = new_title
                for r in runs[1:]:
                    r.text = ""
            break
    return s


def move_slide(prs, from_idx, to_idx):
    lst = prs.slides._sldIdLst
    items = list(lst)
    el = items[from_idx]
    lst.remove(el)
    lst.insert(to_idx, el)


prs = Presentation(DECK)

# ── 1) De-cram slide 5: drop the WL metrics table + cross-village box ─────────
s5 = prs.slides[4]
drop = {510,511,512,513,514,515,520,521}
for sh in list(s5.shapes):
    if sh.shape_id in drop:
        sh._element.getparent().remove(sh._element)
# retitle slide-5 subtitle and header to reflect DTM focus
for sh in s5.shapes:
    if sh.shape_id == 500 and sh.has_text_frame:
        for r in [r for p in sh.text_frame.paragraphs for r in p.runs]:
            r.text = ""
        sh.text_frame.paragraphs[0].runs and None
        # rewrite subtitle text
        p = sh.text_frame.paragraphs[0]
        if not p.runs:
            rr = p.add_run()
        rr = p.runs[0]; rr.text = "DTM vertical accuracy — ASPRS leave-out CV vs withheld LiDAR ground returns"
# add a 3-callout DTM highlights strip in the freed space (y ~3.15–4.35)
hl = [("0.5 m COG DTM","Cloud-Optimized GeoTIFF · OGC-compliant · 8 terrain derivatives",NAVY),
      ("Zero vertical bias","Mean error ≈ 0 m — no systematic drift across 4 villages",TEAL),
      ("RMSE 0.10–0.25 m","Scales with LiDAR point density, not algorithm quality",AMBER)]
cw = 3.05
for i,(t,d,c) in enumerate(hl):
    x = 0.25 + i*(cw+0.12)
    box(s5, x, 3.20, cw, 0.40, [(t,12.5,True,WHITE,PP_ALIGN.CENTER)], fill=c, anchor=MSO_ANCHOR.MIDDLE)
    box(s5, x, 3.60, cw, 0.78, [(d,10.5,False,DARK,PP_ALIGN.LEFT)], fill=LT, line=c)

# ── 2) NEW slide: Waterlogging Risk Maps ─────────────────────────────────────
s6 = dup_results_slide(prs, 4, "Waterlogging Risk — One Model, 4 Villages")
box(s6, 0.25, 0.92, 9.5, 0.30,
    [("Same XGBoost model, zero retraining across Gujarat ↔ Punjab — physics-derived labels, no flood records needed",
      10.5,False,GREY,PP_ALIGN.CENTER)])
pic_fit(s6, FIG/"risk_grid.png", 0.30, 1.28, 5.55, 3.78)
box(s6, 0.30, 5.02, 5.55, 0.26,
    [("Red = high risk · model AUC 0.999 in every village",9.5,False,GREY,PP_ALIGN.CENTER)])

# right panel — model scores + transfer + TWI
rx = 6.05; rw = 3.7
box(s6, rx, 1.28, rw, 0.34, [("Model Accuracy (XGBoost, 5-fold CV)",12,True,WHITE,PP_ALIGN.CENTER)],
    fill=TEAL, anchor=MSO_ANCHOR.MIDDLE)
score_lines = [("Village        ROC-AUC   F1",10,True,NAVY,PP_ALIGN.LEFT)]
for k in VK:
    w = MET["villages"][k]["waterlogging"]
    score_lines.append((f"{VLAB[k]:<11} {w['roc_auc']:.4f}   {w['f1']:.3f}",10,False,DARK,PP_ALIGN.LEFT))
box(s6, rx, 1.62, rw, 1.20, score_lines, fill=LT, line=TEAL)

box(s6, rx, 2.95, rw, 0.34, [("Cross-Village Transfer",12,True,WHITE,PP_ALIGN.CENTER)],
    fill=AMBER, anchor=MSO_ANCHOR.MIDDLE)
box(s6, rx, 3.29, rw, 0.86,
    [("Train on village A → test on B",10,False,GREY,PP_ALIGN.CENTER),
     ("Min AUC 0.989   ·   Max AUC 1.000",12.5,True,TEAL,PP_ALIGN.CENTER),
     ("generalises without retraining",10,False,DARK,PP_ALIGN.CENTER)], fill=LT, line=AMBER)

box(s6, rx, 4.28, rw, 0.34, [("Physics Check — TWI HIGH/LOW = 2.1–2.6×",10.5,True,NAVY,PP_ALIGN.CENTER)],
    fill=LT, line=NAVY, anchor=MSO_ANCHOR.MIDDLE)

# ── 3) NEW slide: Drainage Network Design ────────────────────────────────────
s7 = dup_results_slide(prs, 4, "Drainage Network — Costed & Sized")
box(s7, 0.25, 0.92, 9.5, 0.30,
    [("Flow-accumulation routing (MST) + Manning's trapezoidal sizing for a 10-year storm — every channel costed at CPWD DSR 2023-24 rates",
      10.5,False,GREY,PP_ALIGN.CENTER)])
pic_fit(s7, FIG/"result_DEVDI.png", 0.30, 1.28, 4.85, 3.78)
box(s7, 0.30, 5.02, 4.85, 0.26,
    [("Devdi — high-risk zones (red) drained by sized channels (gold/cyan)",9.5,False,GREY,PP_ALIGN.CENTER)])

# right — cost table across 4 villages
tx=5.35; tw=4.4
box(s7, tx, 1.28, tw, 0.34, [("Costed Drainage Plan — 4 Villages",12,True,WHITE,PP_ALIGN.CENTER)],
    fill=NAVY, anchor=MSO_ANCHOR.MIDDLE)
tot_ch=tot_km=tot_cost=0
rows=[("Village        Chan.   Length   Cost",10,True,NAVY,PP_ALIGN.LEFT)]
for k in VK:
    d=MET["villages"][k]["drainage"]
    ch=d["channel_count"]; km=d["total_length_m"]/1000; cost=d["total_cost_inr_lakhs"]
    tot_ch+=ch; tot_km+=km; tot_cost+=cost
    rows.append((f"{VLAB[k]:<11} {ch:>5}   {km:>5.1f} km  ₹{cost:>5.0f}L",10,False,DARK,PP_ALIGN.LEFT))
rows.append((f"{'TOTAL':<11} {tot_ch:>5}   {tot_km:>5.0f} km  ₹{tot_cost:>5.0f}L",10,True,RED,PP_ALIGN.LEFT))
box(s7, tx, 1.62, tw, 1.55, rows, fill=LT, line=NAVY)

box(s7, tx, 3.30, tw, 0.34, [("Design Parameters",12,True,WHITE,PP_ALIGN.CENTER)],
    fill=TEAL, anchor=MSO_ANCHOR.MIDDLE)
box(s7, tx, 3.64, tw, 1.05,
    [("Manning's n 0.025 · side slope 1.5H:1V · i = 50 mm/hr · C = 0.65",10.5,False,DARK,PP_ALIGN.LEFT),
     ("Velocity 0.30–2.00 m/s · concrete lining only where flow is fast",10.5,False,DARK,PP_ALIGN.LEFT),
     ("Cost: CPWD DSR 2023-24 depth-banded, avg ₹231–267/m incl. 12% GST",10.5,False,DARK,PP_ALIGN.LEFT),
     ("0 capacity exceedances · 100% stream coverage",10.5,True,TEAL,PP_ALIGN.LEFT)],
    fill=LT, line=TEAL)

# ── 4) Reorder: place the two new slides right after Results (idx 4) ──────────
n = len(prs.slides)              # new slides are at indices n-2, n-1
move_slide(prs, n-2, 5)          # WL maps  → position 6
move_slide(prs, n-1, 6)          # Drainage → position 7

prs.save(DECK)
print(f"Saved {DECK.name} — now {len(prs.slides)} slides.")
order = []
for s in prs.slides:
    for sh in s.shapes:
        if sh.is_placeholder and sh.has_text_frame and sh.text_frame.text.strip():
            order.append(sh.text_frame.text[:40]); break
for i,t in enumerate(order): print(f"  {i+1}. {t}")
